"""
Generate docs/service-launch-roadmap.pptx from the launch roadmap content.
Requires: pip install python-pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "service-launch-roadmap.pptx"


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1F, 0x2E)
    p.alignment = PP_ALIGN.LEFT

    sub = tf.add_paragraph()
    sub.text = subtitle
    sub.font.size = Pt(14)
    sub.font.color.rgb = RGBColor(0x4A, 0x55, 0x68)
    sub.space_before = Pt(12)


def _add_bullet_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tt = title_box.text_frame
    tp = tt.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(24)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0x1A, 0x1F, 0x2E)

    body = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(8.7), Inches(5.5))
    bf = body.text_frame
    bf.word_wrap = True
    for i, line in enumerate(lines):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        para.text = line
        para.font.size = Pt(13)
        para.font.color.rgb = RGBColor(0x2D, 0x37, 0x41)
        para.space_after = Pt(6)
        para.level = 0


def _add_table_slide(
    prs: Presentation,
    title: str,
    headers: tuple[str, ...],
    rows: list[tuple[str, ...]],
    col_widths_inches: tuple[float, ...],
) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.65))
    tt = title_box.text_frame
    tp = tt.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(22)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(0x1A, 0x1F, 0x2E)

    ncols = len(headers)
    nrows = 1 + len(rows)
    total_w = sum(col_widths_inches)
    left = Inches(0.45)
    top = Inches(1.05)
    width = Inches(total_w)
    height = Inches(min(0.35 * nrows + 0.2, 4.8))
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x31, 0x41, 0x55)

    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(0x2D, 0x37, 0x41)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "BroGourmet — 최종 서비스 출시 로드맵",
        "체크리스트 · 비즈니스 후보 · 요약  |  기준일: 2026-04-19 (KST)",
    )

    _add_bullet_slide(
        prs,
        "로드맵 흐름 (요약)",
        [
            "【현재】 BroG·지도·MyG·커뮤니티·게임 / 가입·인증·JWT / 결제·KCP·포인트 / 스폰서·이벤트·관리자 / nginx·CORS·Vite",
            "【출시 직전】 법무·개인정보 / 운영·모니터링·백업 / PG·도메인·프로덕션 키 / 신고·모더레이션",
            "【성장·수익】 스폰서·지역 광고 / 가맹·이벤트·유료 노출 / 데이터·추천·B2B",
        ],
    )

    _add_table_slide(
        prs,
        "우선순위 높은 수정·보완",
        ("#", "항목", "세부", "완료"),
        [
            ("2.1", "출시 빌드·환경", "프로덕션 키 분리·www에서 LAN API 금지", "☐"),
            ("2.2", "단계적 오픈", "BROG_ONLY 등 단계·메뉴 합의·문서화", "☐"),
            ("2.3", "법·정책", "약관·개인정보·위치·결제/환불·동의", "☐"),
            ("2.4", "개발용 노출", "가입/인증 DEV 문구 프로덕션 숨김", "☐"),
            ("2.5", "보안", "CSP·입력검증·의존성 (쿠키는 후속 가능)", "☐"),
            ("2.6", "결제·정산", "콜백·취소·영수증·CS", "☐"),
            ("2.7", "운영", "백업·복구·헬스·로그", "☐"),
            ("2.8", "신뢰", "신고·차단·검토 최소 설계", "☐"),
        ],
        (0.55, 1.85, 5.1, 0.6),
    )

    _add_table_slide(
        prs,
        "추가 기능 아이디어",
        ("#", "아이디어", "메모", "완료"),
        [
            ("3.1", "온보딩", "BroG / MyG / 지도 짧은 안내", "☐"),
            ("3.2", "알림", "웹 푸시·알림톡(정책·비용)", "☐"),
            ("3.3", "검색·발견", "즐겨찾기·최근 본 매장", "☐"),
            ("3.4", "가맹/사장님", "조회 등 최소 지표", "☐"),
        ],
        (0.55, 2.0, 5.0, 0.6),
    )

    _add_table_slide(
        prs,
        "비즈니스 모델 후보 (제품 연결)",
        ("축", "설명", "연결"),
        [
            ("스폰서·광고", "구·카테고리·지도 슬롯", "스폰서 스페이스·클릭 측정"),
            ("가맹·이벤트", "등록비·참가비", "merchant intent·Payment"),
            ("포인트", "충전 후 사용", "point_charge·소비처 정의"),
            ("구독", "프리미엄 혜택", "payment tab=user·과금 정의"),
            ("리드", "문의·예약 중개", "상세·외부 링크"),
            ("B2B 리포트", "익명 집계 리포트", "프라이버시 가이드"),
        ],
        (1.35, 3.25, 3.9),
    )

    _add_bullet_slide(
        prs,
        "새 방향 아이디어 (선택)",
        [
            "「이 동네 오늘의 메뉴」 큐레이션 → 스폰서와 묶기",
            "점메추(SADARI) × 스폰서 — 비침습 한 줄",
            "지역 미디어 제휴 — 스폰서 패키지 공동 판매",
        ],
    )

    _add_bullet_slide(
        prs,
        "한 줄 요약",
        [
            "기술적으로는 지도·UGC·결제·스폰서·관리까지 MVP 이상.",
            "최종 서비스까지 남은 일은 법무·운영·실결제·신뢰(신고/모더레이션)·수익 설계를 제품에 녹이는 일.",
            "",
            "문서: docs/service-launch-roadmap.md  |  재생성: python scripts/generate_service_launch_ppt.py",
        ],
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
